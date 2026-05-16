# Ported and enhanced from the original AI Chatbot (g:\Desktop\Uni\2025启航计划\AI Chatbot\deepseek_client.py)
import base64
import io
import os
import re
from datetime import datetime
from typing import Optional, List, Dict, Tuple
from openai import OpenAI
from PyPDF2 import PdfReader
from pptx import Presentation
from docx import Document
from app.core.config import get_settings

settings = get_settings()

client = OpenAI(
    api_key=settings.DEEPSEEK_API_KEY,
    base_url=settings.DEEPSEEK_API_BASE,
)

# ==================== File Parsers ====================

def extract_text_from_pdf(pdf_data: bytes) -> str:
    text = ""
    with io.BytesIO(pdf_data) as f:
        reader = PdfReader(f)
        for page in reader.pages:
            text += page.extract_text() or ""
    return text


def extract_text_from_pptx(pptx_data: bytes) -> str:
    text = ""
    with io.BytesIO(pptx_data) as f:
        prs = Presentation(f)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    return text


def extract_text_from_docx(docx_data: bytes) -> str:
    text = ""
    with io.BytesIO(docx_data) as f:
        doc = Document(f)
        for para in doc.paragraphs:
            text += para.text + "\n"
    return text


# ==================== PDF Page Cache ====================

pdf_page_cache: Dict[str, List[str]] = {}
pdf_metadata_cache: Dict[str, Dict] = {}


def extract_pdf_by_page(pdf_data: bytes) -> List[str]:
    pages = []
    with io.BytesIO(pdf_data) as f:
        reader = PdfReader(f)
        for i, page in enumerate(reader.pages):
            page_text = page.extract_text() or ""
            pages.append(page_text)
    return pages


def get_cached_pdf_pages(course_name: str, chapter_index: int) -> Optional[List[str]]:
    cache_key = f"{course_name}-ch{chapter_index}"
    return pdf_page_cache.get(cache_key)


def cache_pdf_pages(course_name: str, chapter_index: int, pages: List[str]):
    cache_key = f"{course_name}-ch{chapter_index}"
    pdf_page_cache[cache_key] = pages
    pdf_metadata_cache[cache_key] = {
        "total_pages": len(pages),
        "cached_at": datetime.now().isoformat(),
    }


def extract_page_numbers(question: str) -> List[int]:
    page_numbers = []
    chinese_patterns = [
        r'第(\d+)页',
        r'第(\d+)\s*[-~到]\s*(\d+)页',
    ]
    english_patterns = [
        r'page\s*(\d+)',
        r'pages?\s*(\d+)\s*[-~to]+\s*(\d+)',
    ]
    for pattern in chinese_patterns + english_patterns:
        matches = re.finditer(pattern, question, re.IGNORECASE)
        for match in matches:
            groups = match.groups()
            if len(groups) >= 1:
                page_numbers.append(int(groups[0]))
            if len(groups) >= 2 and groups[1] is not None:
                start = int(groups[0])
                end = int(groups[1])
                page_numbers.extend(range(start, end + 1))
    return sorted(set(page_numbers))


def select_relevant_pages(question: str, all_pages: List[str], include_page_numbers: bool = True) -> str:
    page_numbers = extract_page_numbers(question)
    if page_numbers:
        selected_pages = []
        for page_num in page_numbers:
            if 1 <= page_num <= len(all_pages):
                if include_page_numbers:
                    selected_pages.append(f"[Page {page_num}]\n{all_pages[page_num - 1]}")
                else:
                    selected_pages.append(all_pages[page_num - 1])
        if selected_pages:
            return "\n\n".join(selected_pages)

    important_keywords = [
        '概念', '定义', '公式', '定理', '原理', '性质', '特点', '特征',
        '例子', '示例', '例题', '应用', '用途', '方法', '步骤', '流程',
        '总结', '小结', '要点', '重点', '难点', '关键', '结论',
        'concept', 'definition', 'formula', 'theorem', 'principle', 'property',
        'characteristic', 'feature', 'example', 'sample', 'application', 'use',
        'method', 'step', 'procedure', 'process', 'summary', 'key point',
        'conclusion', 'important', 'critical',
    ]

    relevant_pages = []
    for i, page_content in enumerate(all_pages):
        page_lower = page_content.lower()
        for keyword in important_keywords:
            if keyword.lower() in question.lower() and keyword.lower() in page_lower:
                if include_page_numbers:
                    relevant_pages.append(f"[Page {i+1}]\n{page_content}")
                else:
                    relevant_pages.append(page_content)
                break
        if len(relevant_pages) >= 3:
            break

    if relevant_pages:
        return "\n\n".join(relevant_pages)

    default_pages = all_pages[:3]
    if include_page_numbers:
        return "\n\n".join([f"[Page {i+1}]\n{content}" for i, content in enumerate(default_pages)])
    return "\n\n".join(default_pages)


# ==================== System Prompts ====================

SYSTEM_PROMPTS = {
    "chatbot": "You are a general-purpose chatbot, friendly and helpful. Please use Markdown for formatting. Default to responding in English unless specifically asked to use another language.",
    "physics": "You are an expert physicist and educator. Answer clearly, use Markdown, LaTeX for formulas.",
    "math": "You are a patient math tutor. Help users step-by-step using Markdown and LaTeX.",
    "circuit": "You are an experienced circuit engineer. Answer with professional terminology using Markdown and LaTeX.",
}

VOICE_SYSTEM_PROMPT = (
    "You are a conversational voice assistant. Reply in natural, colloquial language suitable for speech. "
    "Do NOT use Markdown, headings, lists, code blocks, inline code markers, or raw URLs. "
    "Avoid special symbols like '#', '*', '`'. Keep answers concise and well-structured as plain sentences. "
    "If the user asked in Chinese, reply in Chinese; if in English, reply in English."
)


# ==================== AI Response ====================

async def get_deepseek_response(
    user_question: str,
    mode: str = "chatbot",
    file_data: Optional[bytes] = None,
    filename: Optional[str] = None,
    conversation_history: Optional[List[Dict[str, str]]] = None,
) -> Dict:
    """
    Returns dict with keys: response (str), context_source (str), filename_display (str)
    """
    if not settings.DEEPSEEK_API_KEY:
        return {"response": "Error: DeepSeek API key is not configured.", "context_source": "none", "filename_display": ""}

    try:
        messages = []
        model_name = settings.DEEPSEEK_MODEL_NAME
        file_context = ""
        context_source = "none"
        filename_display = ""
        course_name = ""
        chapter_index = -1

        # --------- File Upload ---------
        if file_data and filename:
            is_image = filename.lower().endswith(('.png', '.jpg', '.jpeg', '.gif', '.webp'))
            if is_image:
                model_name = "deepseek-vl-chat"
                base64_image = base64.b64encode(file_data).decode('utf-8')
                messages = [{"role": "user", "content": [
                    {"type": "text", "text": user_question or "Please describe this image"},
                    {"type": "image_url", "image_url": {"url": f"data:image/jpeg;base64,{base64_image}"}}
                ]}]
            else:
                if filename.lower().endswith('.pdf'):
                    pages = extract_pdf_by_page(file_data)
                    file_context = select_relevant_pages(user_question, pages, include_page_numbers=True)
                elif filename.lower().endswith('.pptx'):
                    file_context = extract_text_from_pptx(file_data)
                elif filename.lower().endswith('.docx'):
                    file_context = extract_text_from_docx(file_data)
                else:
                    return {"response": f"Unsupported file type: '{filename}'", "context_source": "none", "filename_display": ""}
                context_source = "manual_file"
                filename_display = filename

        # --------- Chapter Mode ---------
        elif '-ch' in mode:
            try:
                course_name_slug, chapter_index_str = mode.split('-ch')
                course_name = course_name_slug.replace('-', ' ')
                chapter_index = int(chapter_index_str)

                script_dir = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
                base_path = os.path.join(script_dir, "knowledge_base", course_name)
                txt_path = os.path.join(base_path, f"chapter_{chapter_index + 1}.txt")
                pdf_path = os.path.join(base_path, f"chapter_{chapter_index + 1}.pdf")

                if os.path.exists(txt_path):
                    with open(txt_path, 'r', encoding='utf-8') as f:
                        file_context = f.read()
                    context_source = "txt"
                    filename_display = f"Chapter {chapter_index + 1} Knowledge Base"
                elif os.path.exists(pdf_path):
                    cached_pages = get_cached_pdf_pages(course_name, chapter_index)
                    if cached_pages is None:
                        with open(pdf_path, 'rb') as f:
                            pdf_bytes = f.read()
                            pages = extract_pdf_by_page(pdf_bytes)
                            cache_pdf_pages(course_name, chapter_index, pages)
                    else:
                        pages = cached_pages
                    file_context = select_relevant_pages(user_question, pages, include_page_numbers=True)
                    context_source = "pdf"
                    filename_display = f"Chapter {chapter_index + 1} PDF"
                else:
                    file_context = ""
                    context_source = "none"
            except (ValueError, IndexError):
                file_context = ""
                context_source = "none"

        # --------- Build System Prompt ---------
        if file_context:
            if context_source in ["txt", "pdf"]:
                page_query_hint = ""
                if extract_page_numbers(user_question):
                    page_query_hint = "\nNote: The user is querying specific page content. Please accurately reference the [Page X] content."

                system_prompt = (
                    "You are an expert teaching assistant in the field of Electronic and Information Engineering. "
                    "Answer questions using the following priority:\n"
                    "1. First, answer based on the provided context (TXT knowledge base or PDF content)\n"
                    "2. If the context does not contain the answer, answer based on your expertise (restricted to Electronic and Information Engineering)\n"
                    "3. If the user asks about a specific page, accurately reference the [Page X] content\n"
                    "4. Do not answer questions unrelated to Electronic and Information Engineering\n"
                    f"{page_query_hint}\n"
                    "Respond in the same language (Chinese or English) as the user."
                )
                combined_input = (
                    f"**Context ({filename_display}):**\n---\n{file_context}\n---\n\n"
                    f"**User Question:**\n{user_question}"
                )
            else:
                system_prompt = SYSTEM_PROMPTS.get(mode, SYSTEM_PROMPTS["chatbot"])
                combined_input = user_question
        else:
            system_prompt = SYSTEM_PROMPTS.get(mode, SYSTEM_PROMPTS["chatbot"])
            combined_input = user_question

        if mode == 'voice_chat' or mode.startswith('voice'):
            if system_prompt:
                system_prompt = system_prompt + "\n\n" + VOICE_SYSTEM_PROMPT
            else:
                system_prompt = VOICE_SYSTEM_PROMPT

        # --------- Build Messages ---------
        if not messages:
            messages = [{"role": "system", "content": system_prompt}]
            if conversation_history and len(conversation_history) > 0:
                for msg in conversation_history:
                    if msg.get("role") in ["user", "assistant"]:
                        messages.append({"role": msg["role"], "content": msg["content"]})
            messages.append({"role": "user", "content": combined_input})

        # --------- Call API ---------
        response = client.chat.completions.create(
            model=model_name,
            messages=messages,
            stream=False,
        )
        ai_response = response.choices[0].message.content.strip()

        return {
            "response": ai_response,
            "context_source": context_source,
            "filename_display": filename_display,
        }

    except Exception as e:
        error_message = f"Error calling DeepSeek API: {str(e)}"
        print(error_message)
        return {"response": error_message, "context_source": "none", "filename_display": ""}
